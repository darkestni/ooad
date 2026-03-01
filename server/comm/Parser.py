from pathlib import Path
import re
from collections import Counter
import io
import zipfile

# 首选使用 PyPDF2 / python-docx / python-pptx
try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

# 可选 OCR 依赖
try:
    import pytesseract
    from pdf2image import convert_from_path
    from PIL import Image, ImageFilter, ImageOps

    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False


class ParserError(Exception):
    pass


class Parser:
    """返回结构化解析结果的解析器（支持深度 OCR）。"""

    IMAGE_EXTS = {'.png', '.jpg', '.jpeg', '.tiff', '.tif', '.bmp', '.webp'}

    def __init__(self):
        pass

    def parse_file(self, file_path, ocr: bool = True, ocr_lang: str = "chi_sim+eng"):
        """
        解析文件，返回结构化结果。
        - 仅输出纯文本信息（包含图片的 OCR 文本），不保留/返回任何二进制图片数据；
          上层再做 chunk → 写入向量库时，自然只会存 file_id / user / chunk[]。
        - ocr_lang 直接透传给 pytesseract，例如:
            - "chi_sim"
            - "eng"
            - "chi_sim+eng"（推荐）
        - ocr=True 时，不仅会对纯图页面 OCR，还会尝试提取嵌入图片进行 OCR。
        """
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"{file_path} 不存在")

        meta = {'source': str(file_path)}
        ext = file_path.suffix.lower()
        meta['ext'] = ext
        pages = []

        # 检查 OCR 依赖
        if ocr and not OCR_AVAILABLE:
            print("Warning: OCR 已启用但缺少依赖 (pytesseract/PIL/pdf2image)，将跳过图片识别。")
            ocr = False

        if ext == '.pdf':
            if PdfReader is None:
                raise ParserError('缺少 PyPDF2，无法解析 PDF。')
            pages = self._parse_pdf(file_path, ocr=ocr, ocr_lang=ocr_lang)

        elif ext == '.docx':
            if Document is None:
                raise ParserError('缺少 python-docx。')
            # Word 包含两部分：文档流文本 + 嵌入图片文本
            text_content = self._parse_word_text(file_path)
            image_content = self._parse_docx_images(file_path, ocr, ocr_lang) if ocr else ""
            # 合并为一个“长页”
            combined = (text_content or "") + ("\n" + image_content if image_content else "")
            pages = [combined]

        elif ext == '.pptx':
            if Presentation is None:
                raise ParserError('缺少 python-pptx。')
            pages = self._parse_ppt(file_path, ocr=ocr, ocr_lang=ocr_lang)

        elif ext in ('.md', '.markdown'):
            pages = [self._parse_markdown(file_path, ocr=ocr, ocr_lang=ocr_lang)]

        elif ext in ('.doc', '.ppt'):
            raise ParserError("不支持二进制 .doc / .ppt，请转换为 .docx / .pptx。")

        elif ext in self.IMAGE_EXTS:
            # 纯图片文件：直接 OCR 成一个“页面”
            pages = [self._parse_image(file_path, ocr=ocr, ocr_lang=ocr_lang)]

        else:
            raise ParserError(f"不支持的文件格式: {ext}")

        cleaned_pages = self._clean_pages(pages)
        text = "\n\n".join(cleaned_pages)

        return {
            'text': text,
            'pages': cleaned_pages,
            'meta': meta,
        }

    # ========================== 核心 OCR 辅助 ==========================

    def _preprocess_image_for_ocr(self, img):
        """图像预处理：转灰度、放大、增强对比度。"""
        try:
            # 统一转 RGB 再处理
            if img.mode != 'RGB':
                img = img.convert('RGB')

            # 针对小图简单放大一倍，避免字太小难以识别
            w, h = img.size
            if max(w, h) < 1000:
                img = img.resize((w * 2, h * 2), Image.LANCZOS)

            # 转灰度
            img = ImageOps.grayscale(img)
            # 轻微中值滤波去噪
            img = img.filter(ImageFilter.MedianFilter(size=3))
            # 自动对比度
            img = ImageOps.autocontrast(img)

            return img
        except Exception:
            return img

    def _ocr_image_bytes(self, image_data, lang: str = "chi_sim+eng"):
        """将内存中的二进制图片数据转换为文本，仅返回识别出的字符串。"""
        if not OCR_AVAILABLE or not image_data:
            return ""
        try:
            with io.BytesIO(image_data) as f:
                img = Image.open(f)
                img = self._preprocess_image_for_ocr(img)
                text = pytesseract.image_to_string(img, lang=lang)
                return text.strip()
        except Exception:
            # print(f"OCR Error: {e}")
            return ""

    # ========================== 各格式解析器 ==========================

    def _parse_pdf(self, file_path: Path, ocr: bool = True, ocr_lang: str = "chi_sim+eng"):
        """
        解析 PDF：
        - 优先使用 PdfReader 提取文本
        - 同时对 page.images 做 OCR（如果可用）
        - 若整页几乎无文本，则渲染整页为图片做兜底 OCR
        """
        reader = PdfReader(str(file_path))
        pages_text = []

        for i, page in enumerate(reader.pages, start=1):
            parts = []

            # 1. 提取可选文本
            try:
                text = page.extract_text() or ""
                if text.strip():
                    parts.append(text)
            except Exception:
                pass

            # 2. 提取页面内嵌图片 (PyPDF2 >= 3.0.0 支持 page.images)
            if ocr and OCR_AVAILABLE:
                try:
                    if hasattr(page, 'images') and page.images:
                        img_texts = []
                        for img_file_obj in page.images:
                            extracted = self._ocr_image_bytes(img_file_obj.data, lang=ocr_lang)
                            if extracted:
                                img_texts.append(f"[Image Text]: {extracted}")
                        if img_texts:
                            parts.append("\n".join(img_texts))
                except Exception:
                    pass

            # 3. 兜底：如果整页几乎空白，尝试渲染整页为图片 OCR
            full_page_text = "\n".join(parts)
            if len(full_page_text.strip()) < 10 and ocr and OCR_AVAILABLE:
                try:
                    pil_pages = convert_from_path(str(file_path), first_page=i, last_page=i)
                    if pil_pages:
                        fallback_img = self._preprocess_image_for_ocr(pil_pages[0])
                        fallback_text = pytesseract.image_to_string(fallback_img, lang=ocr_lang)
                        full_page_text = fallback_text.strip()
                except Exception:
                    pass

            header = f"---PAGE {i}---\n"
            pages_text.append(header + full_page_text.strip())

        return pages_text

    def _parse_word_text(self, file_path: Path) -> str:
        """仅提取 Word 文档流中的文本。"""
        doc = Document(str(file_path))
        return "\n".join([p.text for p in doc.paragraphs])

    def _parse_docx_images(self, file_path: Path, ocr: bool, ocr_lang: str) -> str:
        """
        从 .docx (zip) 包的 word/media 目录中直接提取所有图片。
        缺点：无法确定图片在文中的确切位置。
        策略：将所有图片识别结果附在文档末尾。
        """
        if not ocr or not OCR_AVAILABLE:
            return ""

        extracted_texts = []
        try:
            with zipfile.ZipFile(file_path) as z:
                # 找到所有媒体文件
                media_files = [f for f in z.namelist() if f.startswith('word/media/')]
                for media in media_files:
                    # 简单判断是否为图片
                    if any(media.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.bmp']):
                        img_data = z.read(media)
                        text = self._ocr_image_bytes(img_data, lang=ocr_lang)
                        if text:
                            extracted_texts.append(
                                f"---Embedded Image ({Path(media).name})---\n{text}"
                            )
        except Exception as e:
            print(f"Docx Image Extraction failed: {e}")

        return "\n".join(extracted_texts)

    def _parse_ppt(self, file_path: Path, ocr: bool = True, ocr_lang: str = "chi_sim+eng"):
        """
        解析 PPTX：
        - 提取每个 slide 的文字（文本框 + 表格）
        - 对每个包含 image 的 shape 做 OCR，并附上 [Slide Image]: 前缀
        """
        prs = Presentation(str(file_path))
        slides = []
        for si, slide in enumerate(prs.slides, start=1):
            texts = []

            # 遍历所有 Shape
            for shape in slide.shapes:
                # 1. 文本框 / 表格
                if hasattr(shape, 'text') and shape.text.strip():
                    texts.append(shape.text.strip())

                if hasattr(shape, 'has_table') and shape.has_table:
                    try:
                        for r in shape.table.rows:
                            texts.append("\t".join(cell.text.strip() for cell in r.cells))
                    except Exception:
                        pass

                # 2. 图片 Shape (Type 13 is PICTURE)
                # 使用 hasattr(shape, 'image') 判断更通用
                if ocr and OCR_AVAILABLE and hasattr(shape, 'image'):
                    try:
                        img_text = self._ocr_image_bytes(shape.image.blob, lang=ocr_lang)
                        if img_text:
                            texts.append(f"[Slide Image]: {img_text}")
                    except Exception:
                        pass

            slide_content = "\n".join(texts)
            slides.append(f"---SLIDE {si}---\n" + slide_content)
        return slides

    def _parse_markdown(self, file_path: Path, ocr: bool = True, ocr_lang: str = "chi_sim+eng") -> str:
        """
        解析 Markdown：
        - 读取纯文本
        - 扫描本地图片链接 (![]() 形式)，对能找到的本地图片做 OCR，并在原图后面追加 [Image OCR]: 文本
        """
        # 1. 读取基础文本
        content = ""
        for enc in ("utf-8", "utf-8-sig", "gbk", "latin-1"):
            try:
                with open(file_path, 'r', encoding=enc) as f:
                    content = f.read()
                break
            except UnicodeDecodeError:
                continue

        if not content:
            return ""

        # 2. 扫描本地图片链接并 OCR
        if ocr and OCR_AVAILABLE:
            base_dir = file_path.parent

            # 简单的 Markdown 图片正则: ![alt](path)
            img_pattern = re.compile(r'!\[.*?\]\((.*?)\)')

            def replace_with_ocr(match):
                rel_path = match.group(1).split()[0]  # 有时会有 title 属性
                img_path = (base_dir / rel_path).resolve()

                if img_path.exists() and img_path.suffix.lower() in self.IMAGE_EXTS:
                    try:
                        with open(img_path, 'rb') as f:
                            text = self._ocr_image_bytes(f.read(), lang=ocr_lang)
                        if text:
                            return f"{match.group(0)}\n[Image OCR]: {text}\n"
                    except Exception:
                        pass
                return match.group(0)  # 保持原样

            content = img_pattern.sub(replace_with_ocr, content)

        return content

    def _parse_image(self, file_path: Path, ocr: bool = True, ocr_lang: str = "chi_sim+eng") -> str:
        """
        解析单独图片文件：
        - 不返回任何图片二进制，只返回图片 OCR 结果文本。
        """
        if not OCR_AVAILABLE:
            raise ParserError('OCR 依赖未安装。')
        if not ocr:
            return f"---IMAGE {file_path.name}---\n"

        try:
            img = Image.open(str(file_path))
            proc = self._preprocess_image_for_ocr(img)
            text = pytesseract.image_to_string(proc, lang=ocr_lang) or ""
            return f"---IMAGE {file_path.name}---\n{text.strip()}"
        except Exception as e:
            raise ParserError(f"Image Parse Error: {e}")

    def _clean_pages(self, pages):
        """清洗页眉页脚（保持原逻辑）。"""
        pages_lines = []
        for p in pages:
            lines = [ln.strip() for ln in p.splitlines() if ln and ln.strip()]
            pages_lines.append(lines)

        line_page_count = Counter()
        for lines in pages_lines:
            unique = set(lines)
            for ln in unique:
                line_page_count[ln] += 1

        num_pages = max(1, len(pages_lines))
        # 出现频率超过 60% 的行视为页眉页脚
        remove_lines = {ln for ln, cnt in line_page_count.items() if (cnt / num_pages) > 0.6}

        cleaned_pages = []
        for lines in pages_lines:
            filtered = [ln for ln in lines if ln not in remove_lines]
            # 如果一页被清空，则保留原内容，避免丢失信息
            if not filtered and lines:
                filtered = lines
            cleaned_pages.append("\n".join(filtered))

        return cleaned_pages
