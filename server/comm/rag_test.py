# server/comm/rag_test.py
# pytest: skip-file
"""
RAG 全流程测试工具（升级版 + Parser + 多用户 + 图片 OCR 测试版）

测试范围：
1. 自动生成 PDF / DOCX / PPTX / Markdown（含文本 + 图片）并上传
2. OCR 测试：
   - 单张 PNG
   - PDF 内图片
   - DOCX 内图片
   - PPTX 内图片
   - 纯图片作为“教材”上传
3. 向量库写入：
   - 检查每种文件类型能正确切 chunk
   - 统计 chunks 数量
4. 多用户隔离：
   - 同一教材只对对应 user_id 可见
   - 其它用户检索不到这些 chunks
5. Router：
   - 闲聊（不走检索）
   - 教材问题（必须走检索）
6. Retrieval：
   - k=50 召回
   - 词面重叠 + 知识盲区判断
7. GPT 生成回答（需要你本地有可用代理 & 密钥）
8. 文字 + 图片混合提问（image_paths）
9. 真实论文 TWIST.pdf 端到端 RAG 测试

运行方式：
    python server/comm/rag_test.py
"""

import os
import time

from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader

from PIL import Image, ImageDraw

# DOCX
try:
    import docx  # type: ignore

    HAS_DOCX = True
except Exception:
    HAS_DOCX = False

# PPTX
try:
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches

    HAS_PPTX = True
except Exception:
    HAS_PPTX = False

# 导入 RAG 服务（单例）
from server.comm.rag_service import (
    kb_manager,
    assistant,
    OCRProcessorV2,
    get_vectorstore_stats,
)

# ============================================================
# 工具函数
# ============================================================


def now():
    return time.time()


def t_diff(t0):
    return f"{time.time() - t0:.2f}s"


def banner(title):
    print("\n" + "=" * 60)
    print(f"📌 {title}")
    print("=" * 60 + "\n")


# 测试使用的两个用户，用来检测多用户隔离
USER_ID_MAIN = "USER_TEST_MAIN"
USER_ID_OTHER = "USER_TEST_OTHER"
# 新增：单独给真实论文 TWIST.pdf 用的用户 ID，方便和课程教材隔离
USER_ID_TWIST = "USER_TWIST"

# ============================================================
# 基础文档文本
# ============================================================
BASE_PARAGRAPHS = [
    "This is a test document for the Intelligent Learning Assistant project.",
    "Project A is a secret web application project for the OOAD course.",
    "The deadline for Project A is October 27th, and late submissions may be penalized.",
    "The project requires Python, Java, and basic machine learning knowledge.",
    "Students need to implement a Retrieval-Augmented Generation (RAG) module.",
    "The RAG module should support document upload, automatic chunk splitting, and vector search.",
    "The assistant must answer questions based only on the uploaded materials, avoiding hallucinations.",
    "In addition, students should consider stopwords, noise filtering, and prompt injection defenses.",
    "The system is expected to handle PDF, DOCX, PPTX and Markdown files in a robust way.",
    "For evaluation, teachers may ask about deadlines, tech stacks, and specific project requirements.",
]

SECTION_REPEAT = 4  # 文档长度倍率

# ============================================================
# Step 1 — 生成测试图片
# ============================================================
def generate_test_image(path: str):
    """
    生成一张简单的英文测试图：
    - 白底黑字
    - 字体较大，方便 OCR
    """
    img = Image.new("RGB", (600, 200), color=(255, 255, 255))
    draw = ImageDraw.Draw(img)

    # ⚠️ 提高字体大小，增强 OCR 成功率
    try:
        from PIL import ImageFont

        font = ImageFont.truetype("arial.ttf", 28)
    except Exception:
        font = None

    draw.text((30, 80), "This is an OCR Test Image", fill=(0, 0, 0), font=font)
    img.save(path)
    return path


# ============================================================
# Step 1 — 生成 PDF / DOCX / PPTX（含图片）
# ============================================================
def generate_test_pdf(path: str, img: str):
    banner("生成测试 PDF（含图片）")

    c = canvas.Canvas(path)
    y = 820

    for sec in range(SECTION_REPEAT):
        c.drawString(50, y, f"Section {sec + 1}: Detailed Requirements")
        y -= 20

        for line in BASE_PARAGRAPHS:
            c.drawString(50, y, line)
            y -= 20

            if y < 120:
                c.showPage()
                y = 820

        extra = "End-to-end RAG requires document ingestion, chunking, retrieval, and QA."
        c.drawString(50, y, extra)
        y -= 30

        if y < 200:
            c.showPage()
            y = 820

    # 插入 OCR 测试图
    c.drawImage(ImageReader(img), 50, 80, width=300, height=120)
    c.save()
    print(f"[OK] 生成 PDF：{path}")


def generate_test_docx(path: str, img: str):
    if not HAS_DOCX:
        print("[WARN] 未安装 python-docx，跳过 DOCX 生成")
        return False

    banner("生成测试 DOCX（含图片）")

    d = docx.Document()
    d.add_heading("RAG Test DOCX — With OCR Image", level=1)

    for sec in range(SECTION_REPEAT):
        d.add_heading(f"Section {sec + 1}", level=2)
        for line in BASE_PARAGRAPHS:
            d.add_paragraph(line)

        d.add_paragraph("Extra requirement: retrieval quality evaluation.")

    d.add_picture(img, width=docx.shared.Inches(3))
    d.save(path)
    print(f"[OK] 生成 DOCX：{path}")
    return True


def generate_test_pptx(path: str, img: str):
    if not HAS_PPTX:
        print("[WARN] 未安装 python-pptx，跳过 PPTX 生成")
        return False

    banner("生成测试 PPTX（含图片）")

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "RAG Test PPTX — With OCR Image"

    body = slide.placeholders[1].text_frame
    body.text = "Project A Requirements:"

    for line in BASE_PARAGRAPHS:
        p = body.add_paragraph()
        p.text = line

    # 插入图片
    slide.shapes.add_picture(img, Inches(1), Inches(4), width=Inches(3))
    prs.save(path)
    print(f"[OK] 生成 PPTX：{path}")
    return True


def generate_test_markdown(path: str):
    """
    生成一个 Markdown 文件，用于测试 Parser + md 支持。
    """
    banner("生成测试 Markdown 文件")

    lines = [
        "# RAG Test Markdown",
        "",
        "## Project A Overview",
        "",
        "- Course: OOAD Intelligent Learning Assistant",
        "- Deadline: **October 27th**",
        "- Tech Stack: Python, Java, basic ML",
        "",
        "## Requirements",
        "",
        "1. Support PDF / DOCX / PPTX / Markdown uploads.",
        "2. Implement chunking and vector search.",
        "3. Ensure **low hallucination** by strict RAG.",
        "",
        "## Extra Notes",
        "",
        "This Markdown file is used to test the Parser integration.",
    ]
    with open(path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))
    print(f"[OK] 生成 Markdown：{path}")
    return True


# ============================================================
# Step 2 — OCR 测试
# ============================================================
def test_single_image_ocr(path):
    banner("测试：单独 PNG 图片 OCR（OCRProcessorV2）")
    txt = OCRProcessorV2.ocr_image_file(path)
    print(txt)
    return txt


def test_pdf_image_ocr(path):
    banner("测试：PDF 内图片 OCR（原始 OCRProcessor）")
    from server.comm.rag_service import OCRProcessor

    res = OCRProcessor.extract_images_from_pdf(path)
    for r in res:
        print(r)
    return res


def test_docx_image_ocr(path):
    banner("测试：DOCX 内图片 OCR（原始 OCRProcessor）")
    from server.comm.rag_service import OCRProcessor

    res = OCRProcessor.extract_images_from_docx(path)
    for r in res:
        print(r)
    return res


def test_pptx_image_ocr(path):
    banner("测试：PPTX 内图片 OCR（原始 OCRProcessor）")
    from server.comm.rag_service import OCRProcessor

    res = OCRProcessor.extract_images_from_pptx(path)
    for r in res:
        print(r)
    return res


# ============================================================
# Step 3 — 向量库写入测试（带 user_id）
# ============================================================
def upload_and_show(doc_id, path, user_id):
    """
    调用 kb_manager.upload_data：
    - doc_id: 作为 doc_id / file_id
    - path: 文件路径
    - user_id: 当前用户，用于写入 chunk.metadata["user_id"]
    """
    banner(f"向量库写入：doc_id={doc_id}, user={user_id}")

    t0 = now()
    chunks = kb_manager.upload_data(
        doc_id, path, original_name=os.path.basename(path), user_id=user_id
    )
    print(f"[OK] 写入 chunks = {len(chunks)}  耗时 = {t_diff(t0)}")

    print("\n📘 示例 Chunk：")
    for c in chunks[:2]:
        print("------------")
        print(c)

    return chunks


# ============================================================
# Step 4 — RAG 测试（Deadline / TechStack / Gap）
# ============================================================
def test_rag_query_deadline(user_id):
    banner(f"RAG 测试：截止日期（user={user_id}）")
    q = "What is the deadline for Project A?"
    res = assistant.handle_user_query(user_id, q)
    print("retrieval_performed:", res["retrieval_performed"])
    print("matched_chunks:", len(res["matched_chunks"]))
    print("\n💬 答案：")
    print(res["final_answer"])
    return res


def test_rag_query_tech(user_id):
    banner(f"RAG 测试：技术栈（user={user_id}）")
    q = "Which technologies are required for Project A?"
    res = assistant.handle_user_query(user_id, q)
    print("retrieval_performed:", res["retrieval_performed"])
    print("matched_chunks:", len(res["matched_chunks"]))
    print("\n💬 答案：")
    print(res["final_answer"])
    return res


def test_rag_query_gap(user_id):
    banner(f"RAG 测试：知识盲区（量子引力）（user={user_id}）")
    q = "Explain quantum gravity."
    res = assistant.handle_user_query(user_id, q)
    print("retrieval_performed:", res["retrieval_performed"])
    print("matched_chunks:", len(res["matched_chunks"]))
    print("\n💬 答案：")
    print(res["final_answer"])
    return res


# ============================================================
# Step 5 — Router 测试（闲聊 / 教材问题）
# ============================================================
def test_router_small_talk(user_id):
    banner(f"Router 测试：闲聊不应走检索（user={user_id}）")
    q = "你好，最近天气怎么样？顺便介绍一下你自己。"
    res = assistant.handle_user_query(user_id, q)
    print("retrieval_performed:", res["retrieval_performed"])
    print("matched_chunks:", len(res["matched_chunks"]))
    print("\n💬 答案：")
    print(res["final_answer"])
    return res


def test_router_course_question(user_id):
    banner(f"Router 测试：教材问题必须走检索（user={user_id}）")
    q = "根据教材，这张图相关的 Project A 的要求是什么？"
    res = assistant.handle_user_query(user_id, q)
    print("retrieval_performed:", res["retrieval_performed"])
    print("matched_chunks:", len(res["matched_chunks"]))
    print("\n💬 答案：")
    print(res["final_answer"])
    return res


# ============================================================
# Step 6 — 多用户隔离测试
# ============================================================
def test_multi_user_isolation():
    """
    只给 USER_ID_MAIN 上传教材，不给 USER_ID_OTHER 上传。
    然后同一问题各问一次，对比 matched_chunks 数量。
    """
    banner("多用户隔离测试：同一问题，不同用户的检索结果对比")

    q = "What is the deadline for Project A?"

    res_main = assistant.handle_user_query(USER_ID_MAIN, q)
    res_other = assistant.handle_user_query(USER_ID_OTHER, q)

    print(
        f"user={USER_ID_MAIN}: retrieval={res_main['retrieval_performed']}, "
        f"matched_chunks={len(res_main['matched_chunks'])}"
    )
    print(
        f"user={USER_ID_OTHER}: retrieval={res_other['retrieval_performed']}, "
        f"matched_chunks={len(res_other['matched_chunks'])}"
    )

    print("\n💬 USER_ID_MAIN 答案：")
    print(res_main["final_answer"])

    print("\n💬 USER_ID_OTHER 答案：")
    print(res_other["final_answer"])

    return res_main, res_other


# ============================================================
# Step 7 — 文字 + 图片混合提问
# ============================================================
def test_rag_query_with_image(user_id, image_path):
    """
    使用 handle_user_query 的 image_paths 参数，测试多图入口：
    - query 为“请根据图片内容回答：图片里写了什么”
    - image_paths 传入我们生成的 OCR 测试图
    """
    banner(f"RAG 测试：文字 + 图片混合提问（user={user_id}）")

    q = "请根据我上传的图片内容，告诉我图片上用英文写了什么？"
    res = assistant.handle_user_query(user_id, q, image_paths=[image_path])
    print("retrieval_performed:", res["retrieval_performed"])
    print("matched_chunks:", len(res["matched_chunks"]))
    print("\n💬 答案：")
    print(res["final_answer"])
    return res


# ============================================================
# Step 8 — 真实论文 TWIST.pdf RAG 测试
# ============================================================
def test_twist_paper_rag(user_id):
    """
    使用真实论文 TWIST.pdf 进行一次端到端 RAG 测试。

    要求：
    - 论文已通过 kb_manager.upload_data 写入向量库（doc_id="DOC_TWIST"）
    - user_id 为 USER_ID_TWIST
    """
    banner(f"实际论文 RAG 测试：TWIST.pdf（user={user_id}）")

    q = (
        "What does the Related Works part talk about"
        "according to the essay"
    )

    res = assistant.handle_user_query(user_id, q)
    print("retrieval_performed:", res["retrieval_performed"])
    print("matched_chunks:", len(res["matched_chunks"]))

    print("\n💬 答案：")
    print(res["final_answer"])

    # 可选：打印一个示例 chunk，看一下模型到底看的是论文的哪一段
    if res.get("matched_chunks"):
        print("\n📘 示例 Chunk：")
        print("------------")
        print(res["matched_chunks"][0])
        print("------------")

    return res


# ============================================================
# 主入口
# ============================================================
if __name__ == "__main__":

    banner("🚀 Intelligent Learning Assistant — 全流程测试启动")

    IMG = "ocr_test.png"
    PDF = "rag_test.pdf"
    DOCX = "rag_test.docx"
    PPTX = "rag_test.pptx"
    MD = "rag_test.md"
    # 真实论文文件名（请将 TWIST.pdf 放在与本脚本相同目录）
    TWIST = "TWIST.pdf"

    # Step 0：打印向量库初始统计
    banner("向量库初始统计")
    try:
        get_vectorstore_stats()
    except Exception as e:
        print("[WARN] 无法获取向量库统计：", e)

    # Step 1：生成基础 OCR 测试图片
    generate_test_image(IMG)

    # Step 2：生成带图片的 PDF / DOCX / PPTX / Markdown
    generate_test_pdf(PDF, IMG)
    docx_ok = generate_test_docx(DOCX, IMG)
    pptx_ok = generate_test_pptx(PPTX, IMG)
    md_ok = generate_test_markdown(MD)

    # Step 3：测试单图 OCR
    test_single_image_ocr(IMG)

    # Step 4：测试 PDF / DOCX / PPTX 内嵌图片 OCR
    test_pdf_image_ocr(PDF)
    if docx_ok:
        test_docx_image_ocr(DOCX)
    if pptx_ok:
        test_pptx_image_ocr(PPTX)

    # Step 5：上传到向量库（chunk → embedding → chroma），使用同一个 user_id
    upload_and_show("DOC_PDF", PDF, user_id=USER_ID_MAIN)
    if docx_ok:
        upload_and_show("DOC_DOCX", DOCX, user_id=USER_ID_MAIN)
    if pptx_ok:
        upload_and_show("DOC_PPTX", PPTX, user_id=USER_ID_MAIN)
    if md_ok:
        upload_and_show("DOC_MD", MD, user_id=USER_ID_MAIN)

    # 额外：把纯图片当成“教材”上传，测试图片文件作为知识库
    upload_and_show("DOC_IMG_ONLY", IMG, user_id=USER_ID_MAIN)

    # 额外：上传真实论文 TWIST.pdf（如果文件存在）
    if os.path.exists(TWIST):
        upload_and_show("DOC_TWIST", TWIST, user_id=USER_ID_TWIST)
    else:
        print(f"[WARN] 未在当前目录找到 {TWIST}，跳过 TWIST.pdf 上传。")

    # Step 6：RAG 测试（命中 / 不命中）
    test_rag_query_deadline(USER_ID_MAIN)
    test_rag_query_tech(USER_ID_MAIN)
    test_rag_query_gap(USER_ID_MAIN)

    # Step 7：Router 测试（闲聊 / 教材问题）
    test_router_small_talk(USER_ID_MAIN)
    test_router_course_question(USER_ID_MAIN)

    # Step 8：多用户隔离测试（USER_ID_OTHER 没有上传任何教材）
    test_multi_user_isolation()

    # Step 9：文字 + 图片混合提问（image_paths 参数）
    test_rag_query_with_image(USER_ID_MAIN, IMG)

    # Step 10：真实论文 TWIST.pdf RAG 测试（仅在文件存在且已上传时执行）
    if os.path.exists(TWIST):
        test_twist_paper_rag(USER_ID_TWIST)
    else:
        print(f"[WARN] 未找到 {TWIST}，跳过 TWIST.pdf RAG 测试。")

    # Step 11：再次打印向量库统计
    banner("向量库最终统计")
    try:
        get_vectorstore_stats()
    except Exception as e:
        print("[WARN] 无法获取向量库统计：", e)

    banner("🎉 所有测试全部完成！")